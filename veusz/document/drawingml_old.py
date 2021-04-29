import pptx
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Mm
import zipfile
import xml.etree.ElementTree as ET
import os
import shutil
import io
import tempfile

Override = r"{http://schemas.openxmlformats.org/package/2006/content-types}Override"
ChartType = r"application/vnd.openxmlformats-officedocument.drawingml.chart+xml"

def get_pt_val(pt):
    str_value = pt.xpath("./c:v")[0].text #retrieve point value
    value = float(str_value)
    return value

def read_xy(series):
    xVal = {}
    yVal = {}
    ser = series._ser
    x_pts = ser.xpath(".//c:xVal//c:pt") # get all xVals from xml with xpath query
    y_pts = ser.xpath(".//c:yVal//c:pt") # get all yVals from xml with xpath query
    for i in range(len(x_pts)): #loop through all xVals
        x_value = get_pt_val(x_pts[i]) #call function to get each x value
        y_value = get_pt_val(y_pts[i]) #call function to get each y value
        xVal[x_pts[i].idx] = x_value #store x value in dictionary
        yVal[y_pts[i].idx] = y_value # store y value in dictionary

    # in case x & y idx don't have matching pairs return keys that are common to both x & y
    key = set.intersection(*tuple(set(d.keys()) for d in [xVal, yVal]))
    xVal = [xVal[x] for x in key] #create xVal list
    yVal = [yVal[x] for x in key] #create yVal list
    return xVal, yVal

def toMimes(ba):
    print(type(ba))
    tmpdir = tempfile.TemporaryDirectory()
    TEMP_NAME_CLIP = "clip.zip"
    clippath = os.path.join(tmpdir.name, TEMP_NAME_CLIP)
    TEMP_NAME_PPTX = "container"
    pdirpath = os.path.join(tmpdir.name, TEMP_NAME_PPTX)
    pptxpath = os.path.join(tmpdir.name, TEMP_NAME_PPTX + ".pptx")
    pzippath = os.path.join(tmpdir.name, TEMP_NAME_PPTX + ".zip")
    pptxpath_new = os.path.join(tmpdir.name, TEMP_NAME_PPTX + "_new" + ".pptx")

    stream = io.BytesIO(ba)
    # Get DrawingML object from clipboard
    with zipfile.ZipFile(stream, "r") as z:
        with z.open("[Content_Types].xml") as f:
            tree = ET.fromstring(f.read())
        part_names = []
        for child in tree.iterfind(Override):
            content_type = child.attrib["ContentType"]
            if content_type == ChartType:
                part_name = child.attrib["PartName"]
                part_names.append(part_name)
        chart_xmls = []
        for part_name in part_names:
            with io.TextIOWrapper(z.open(part_name.strip("/"), "r")) as f: 
                chart_xml = f.read()
                chart_xmls.append(chart_xml)

    # Create pptx container
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    chart_data = CategoryChartData()
    x, y, cx, cy = Mm(10), Mm(10), Mm(10), Mm(10)
    for i in part_names:
        slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data)
    prs.save(pptxpath)

    # Override
    with zipfile.ZipFile(pptxpath, "r") as z:
        filenames = z.namelist()
        z.extractall(pdirpath)
    for part_name, chart_xml in zip(part_names, chart_xmls):
        with open(os.path.join(pdirpath, "ppt/charts", os.path.basename(part_name)), "w") as f:
            f.write(chart_xml)
    shutil.make_archive(pdirpath, 'zip', pdirpath)
    os.rename(pzippath, pptxpath_new)

    # Re-open
    prs = Presentation(pptxpath_new)
    charts = []
    for i in range(len(part_names)):
        chart = prs.slides[0].shapes[i + 2].chart
        charts.append(chart)

    #graphtexts = []
    data_commands = []

    headers = []
    bodies = []
    graph_num = [str(len(charts))]

    for i, chart in enumerate(charts):         
        
        header = []
        body = []
        
        x_axis = chart.category_axis
        y_axis = chart.value_axis

        header.append("graph")
        #graph_title = chart.chart_title.text_frame.text if chart.has_title else "graph1"
        graph_title = "graph{}".format(str(i+1))
        header.append("'{}'".format(graph_title))
        header.append("'/page1/{}'".format(graph_title))

        body.append("Add('axis', name=u'x', autoadd=False)")
        body.append("Add('axis', name=u'y', autoadd=False)")

        body.append("To(u'x')")
        if x_axis.axis_title.has_text_frame:
            body.append("Set('label', u'{}')".format(x_axis.axis_title.text_frame.text))
        if x_axis.minimum_scale:
            body.append("Set('min', {})".format(x_axis.minimum_scale))
        if x_axis.maximum_scale:
            body.append("Set('max', {})".format(x_axis.maximum_scale))
        body.append("To('..')")

        body.append("To(u'y')")
        body.append("Set('direction', u'vertical')")
        if y_axis.axis_title.has_text_frame:
            body.append("Set('label', u'{}')".format(y_axis.axis_title.text_frame.text))
        if y_axis.minimum_scale:
            body.append("Set('min', {})".format(y_axis.minimum_scale))
        if y_axis.maximum_scale:
            body.append("Set('max', {})".format(y_axis.maximum_scale))
        body.append("To('..')")

        for j, plot in enumerate(chart.plots):
            for k, series in enumerate(plot.series):
                xDataName = series.name if series.name else "xData"
                xDataName += "_{}{}{}".format(i, j, k)
                yDataName = series.name if series.name else "yData"
                yDataName += "_{}{}{}".format(i, j, k)
                plottype = type(plot)
                if plottype is pptx.chart.plot.XyPlot:
                    body.append("Add('xy', name=u'xy{}', autoadd=False)".format(str(k+1)))
                    body.append("To(u'xy{}')".format(str(k+1)))
                    body.append("Set('xData', '{}')".format(xDataName))
                    body.append("Set('yData', '{}')".format(yDataName))
                    body.append("To('..')")

                    xData, yData = read_xy(series)
                    data_commands.append("ImportString(u'`{}`(numeric)','''".format(xDataName))
                    for datapoint in xData:
                        data_commands.append(str(datapoint))
                    data_commands.append("''')")
                    data_commands.append("ImportString(u'`{}`(numeric)','''".format(yDataName))
                    for datapoint in yData:
                        data_commands.append(str(datapoint))
                    data_commands.append("''')")         
                    
                elif plottype is pptx.chart.plot.BarPlot:
                    print("I see, bar")
                elif plottype is pptx.chart.plot.LinePlot:
                    print("Ah, Line plot")
                else:
                    print("what?")

        header.append(str(len(body)))
        headers += header
        bodies += body

    widget_commands = graph_num + headers + bodies
    widget_savetext = '\n'.join(widget_commands) + '\n' + '\n'.join(body) + '\n'
    data_savetext = '\n'.join(data_commands) + '\n'

    tmpdir.cleanup()

    return widget_savetext, data_savetext